"""Génère MAS_Presentation.pptx — 11 slides, mise en page épurée, basée sur docs/."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(r"C:\Users\pc\Desktop\DOSSIERS_BUREAU\MAS-ENSAM-2026")
TEMPLATE = ROOT / "template.pptx"
OUTPUT = ROOT / "MAS_Presentation.pptx"

# Palette (alignée avec la doc Mintlify : vert #0D9373)
GREEN = RGBColor(0x0D, 0x93, 0x73)
GREEN_DARK = RGBColor(0x0A, 0x6E, 0x57)
DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREY_TEXT = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BLUE = RGBColor(0x1A, 0x6B, 0x9A)
ORANGE = RGBColor(0xB4, 0x53, 0x09)
NEUTRAL = RGBColor(0x37, 0x41, 0x51)

prs = Presentation(str(TEMPLATE))
# Vider proprement les slides existantes du template (parts + relations)
def _purge_template_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    slides_part = prs.part
    # Récupérer les rIds des slides référencées
    rIds = [sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            for sldId in list(sldIdLst)]
    # Drop chaque part de slide + sa relation
    for rId in rIds:
        rel = slides_part.rels[rId]
        slide_part = rel.target_part
        # supprimer la part du package
        try:
            del slide_part.package._parts[slide_part.partname]
        except (KeyError, AttributeError):
            pass
        slides_part.drop_rel(rId)
    # Vider l'index XML
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)

_purge_template_slides(prs)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
# On utilise un layout vide pour tout maîtriser ; sinon fallback
for lay in prs.slide_layouts:
    if "Blank" in lay.name or lay.name.lower().startswith("blank"):
        BLANK = lay
        break


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color: RGBColor):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # envoyer en arrière-plan
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    rect = slide.shapes.add_shape(shape, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_accent_bar(slide, y=Inches(0.45), color=GREEN):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.5), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_footer(slide, page, total=11):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "MAS — Multi-Agent Security Scanner", size=10, color=GREY_TEXT)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3),
             f"{page} / {total}", size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, *, subtitle=None):
    add_accent_bar(slide)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(8), Inches(0.35),
             eyebrow, size=11, bold=True, color=GREEN, font="Calibri")
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.8),
             title, size=32, bold=True, color=DARK)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
                 subtitle, size=14, color=GREY_TEXT)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Titre
# ────────────────────────────────────────────────────────────────────────────
s = add_slide()
set_bg(s, DARK)
# halo vert
halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(8), Inches(8))
halo.fill.solid(); halo.fill.fore_color.rgb = GREEN_DARK
halo.line.fill.background(); halo.shadow.inherit = False

add_text(s, Inches(0.8), Inches(2.0), Inches(2.5), Inches(0.4),
         "MAS · ENSAM 2026", size=12, bold=True, color=GREEN)
add_text(s, Inches(0.8), Inches(2.5), Inches(11), Inches(1.6),
         "Multi-Agent Security\nScanner", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.6), Inches(10), Inches(1),
         "Analyse de sécurité de code orchestrée par 8 agents IA spécialisés —\nanalyse statique, raisonnement LLM, patches auto, mémoire persistante.",
         size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

# Tag stack
tags = ["FastAPI", "LangGraph", "Groq · NVIDIA", "Semgrep", "Qdrant", "Rust"]
x = Inches(0.8)
for t in tags:
    w = Inches(0.18 * len(t) + 0.6)
    chip = add_rect(s, x, Inches(6.2), w, Inches(0.4), SLATE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, Inches(6.22), w, Inches(0.36), t, size=11, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
    x += w + Inches(0.12)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problème & Solution
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "01 · CONTEXTE", "Problème & approche",
            subtitle="Pourquoi un pipeline multi-agents pour la sécurité du code.")

col_w = Inches(5.9); col_h = Inches(4.2)
left = Inches(0.6); right = Inches(6.7); top = Inches(2.3)

# Carte problème
add_rect(s, left, top, col_w, col_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
add_text(s, left + Inches(0.4), top + Inches(0.3), col_w - Inches(0.6), Inches(0.4),
         "PROBLÈME", size=11, bold=True, color=ORANGE)
add_text(s, left + Inches(0.4), top + Inches(0.7), col_w - Inches(0.6), Inches(0.6),
         "Outils silos, pas de raisonnement", size=20, bold=True, color=DARK)
bullets_pb = [
    "Semgrep, Bandit, Gosec… chacun isolé",
    "Pas de scoring CVSS contextuel",
    "Failles logiques (IDOR, auth bypass) invisibles",
    "Aucun patch automatique validé",
]
y = top + Inches(1.5)
for b in bullets_pb:
    add_text(s, left + Inches(0.4), y, Inches(0.2), Inches(0.4),
             "•", size=14, color=ORANGE, bold=True)
    add_text(s, left + Inches(0.65), y, col_w - Inches(1), Inches(0.4),
             b, size=13, color=GREY_TEXT)
    y += Inches(0.5)

# Carte solution
add_rect(s, right, top, col_w, col_h, DARK)
add_text(s, right + Inches(0.4), top + Inches(0.3), col_w - Inches(0.6), Inches(0.4),
         "SOLUTION", size=11, bold=True, color=GREEN)
add_text(s, right + Inches(0.4), top + Inches(0.7), col_w - Inches(0.6), Inches(0.6),
         "Pipeline orchestré, 8 agents", size=20, bold=True, color=WHITE)
bullets_sol = [
    "LangGraph route, parallélise, retry",
    "LLM 70B pour le raisonnement sémantique",
    "Scoring CVSS 3.1 + exploitabilité",
    "Patches générés, appliqués, re-scannés",
]
y = top + Inches(1.5)
for b in bullets_sol:
    add_text(s, right + Inches(0.4), y, Inches(0.2), Inches(0.4),
             "▸", size=14, color=GREEN, bold=True)
    add_text(s, right + Inches(0.65), y, col_w - Inches(1), Inches(0.4),
             b, size=13, color=RGBColor(0xCB, 0xD5, 0xE1))
    y += Inches(0.5)

add_footer(s, 2)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Ce que fait la plateforme (7 étapes simplifiées → 4)
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "02 · VUE D'ENSEMBLE", "Du dépôt au rapport, en 4 phases",
            subtitle="Détection → analyse → patch → rapport. Tout est automatisé.")

steps = [
    ("01", "Triage & Scan", "Détection multi-langage, Semgrep + Bandit +\nGosec + SpotBugs + PHPCS en parallèle.", GREEN),
    ("02", "Analyse IA", "Memory safety (Rust) + failles logiques (LLM 70B + RAG Qdrant).", BLUE),
    ("03", "Scoring & Patch", "CVSS 3.1, exploitabilité, unified diff généré\npar LLM si exploitable.", PURPLE),
    ("04", "Validation", "Application du diff, re-scan, retry × 3.\nRapport JSON final.", ORANGE),
]
card_w = Inches(2.95); card_h = Inches(3.6); gap = Inches(0.15)
total_w = card_w * 4 + gap * 3
start_x = (SLIDE_W - total_w) // 2
y = Inches(2.4)
for i, (num, name, desc, col) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    add_rect(s, x, y, card_w, card_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    # Bandeau numéro
    add_rect(s, x, y, card_w, Inches(0.55), col, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, x + Inches(0.3), y + Inches(0.1), card_w - Inches(0.6), Inches(0.4),
             num, size=14, bold=True, color=WHITE)
    add_text(s, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(0.5),
             name, size=18, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), Inches(2),
             desc, size=12, color=GREY_TEXT)

add_footer(s, 3)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Architecture globale
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "03 · ARCHITECTURE", "Vue technique des blocs",
            subtitle="FastAPI → LangGraph → 8 agents → outils & LLM & mémoire.")

# Bandes horizontales
def band(x, y, w, h, title_, items, color):
    add_rect(s, x, y, w, h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, x + Inches(0.3), y + Inches(0.15), w - Inches(0.6), Inches(0.4),
             title_, size=13, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), y + Inches(0.55), w - Inches(0.6), Inches(0.5),
             items, size=11, color=GREY_TEXT)

x0 = Inches(0.6); w_full = Inches(12.13)
band(x0, Inches(2.3), w_full, Inches(0.95),
     "CLIENTS  ·  HTTP REST  +  MCP JSON-RPC 2.0",
     "Entrées du système — 20+ endpoints REST · serveur MCP compatible assistants IA.",
     NEUTRAL)

band(x0, Inches(3.35), Inches(5.95), Inches(0.95),
     "FastAPI  ·  src/api.py",
     "Endpoints, background tasks, SQLite (.memory_cache.db) : users · projets · scans.",
     GREEN)
band(x0 + Inches(6.18), Inches(3.35), Inches(5.95), Inches(0.95),
     "LangGraph Workflow  ·  src/graph/",
     "AgentState partagé, routage conditionnel, parallélisation, retry loop.",
     PURPLE)

band(x0, Inches(4.4), w_full, Inches(0.95),
     "8 Agents  ·  src/agents/",
     "Triage · Scanner · MemorySafety · SemanticAnalyst · ExploitScorer · Patcher · Validator · Report",
     ORANGE)

band(x0, Inches(5.45), Inches(3.93), Inches(1.4),
     "Outils sécurité",
     "Semgrep · Bandit · Gosec ·\nSpotBugs · PHPCS",
     BLUE)
band(x0 + Inches(4.13), Inches(5.45), Inches(3.93), Inches(1.4),
     "LLM Client  ·  src/llm/",
     "Groq Llama-3.1-8B (fast)\nNVIDIA Llama-3.1-70B (strong)",
     PURPLE)
band(x0 + Inches(8.26), Inches(5.45), Inches(3.87), Inches(1.4),
     "Mémoire  ·  src/memory/",
     "Qdrant (vuln_patterns, patches)\nSQLite (users · scans)",
     NEUTRAL)

add_footer(s, 4)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Les 8 agents
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "04 · AGENTS", "Pipeline des 8 agents spécialisés",
            subtitle="Chacun avec un rôle, une entrée et une sortie typées dans AgentState.")

agents = [
    ("1", "Triage", "Détection langages,\nScanTarget[]", GREEN),
    ("2", "Scanner", "Analyse statique\nmulti-outil + cache", GREEN),
    ("3", "MemorySafety", "Buffer overflow,\nUAF — engine Rust", BLUE),
    ("4", "Semantic", "Failles logiques\nLLM 70B + RAG", BLUE),
    ("5", "ExploitScorer", "CVSS 3.1 +\nexploitabilité — LLM 8B", PURPLE),
    ("6", "Patcher", "Unified diffs\nLLM 70B + RAG patches", ORANGE),
    ("7", "Validator", "Patch + re-scan,\nrégressions", ORANGE),
    ("8", "Report", "JSON final,\nstats, top CWEs", NEUTRAL),
]

cols, rows = 4, 2
card_w = Inches(2.95); card_h = Inches(1.85); gap = Inches(0.15)
total_w = card_w * cols + gap * (cols - 1)
start_x = (SLIDE_W - total_w) // 2
start_y = Inches(2.4)

for i, (num, name, desc, color) in enumerate(agents):
    r, c = divmod(i, cols)
    x = start_x + c * (card_w + gap)
    y = start_y + r * (card_h + gap)
    add_rect(s, x, y, card_w, card_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    # Pastille numéro
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.25),
                                Inches(0.55), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background(); badge.shadow.inherit = False
    add_text(s, x + Inches(0.25), y + Inches(0.27), Inches(0.55), Inches(0.55),
             num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.95), y + Inches(0.3), card_w - Inches(1.1), Inches(0.5),
             name, size=15, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.95), card_w - Inches(0.5), Inches(0.9),
             desc, size=11, color=GREY_TEXT)

add_footer(s, 5)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Workflow LangGraph
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, DARK)
add_accent_bar(s, color=GREEN)
add_text(s, Inches(0.6), Inches(0.55), Inches(8), Inches(0.35),
         "05 · WORKFLOW", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.8),
         "Graphe LangGraph & boucle de retry", size=32, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
         "Routage conditionnel, parallélisation MemorySafety ∥ Semantic, retry patching × 3.",
         size=14, color=RGBColor(0xCB, 0xD5, 0xE1))

# Diagramme schématique
def node(x, y, w, h, label, color):
    add_rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, y, w, h, label, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def arrow(x1, y1, x2, y2, color=WHITE):
    line = s.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line

nh = Inches(0.7); nw = Inches(1.6)
y_main = Inches(3.2)
xs = [Inches(0.6), Inches(2.4), Inches(4.2), Inches(6.0), Inches(7.8), Inches(9.6), Inches(11.4)]
# Ligne 1
node(xs[0], y_main, nw, nh, "Triage", GREEN_DARK)
node(xs[1], y_main, nw, nh, "Scanner", GREEN_DARK)
# Parallèle
node(xs[2], Inches(2.55), nw, nh, "Memory\nSafety", BLUE)
node(xs[2], Inches(3.85), nw, nh, "Semantic", BLUE)
node(xs[3], y_main, nw, nh, "Exploit\nScorer", PURPLE)
node(xs[4], y_main, nw, nh, "Patcher", ORANGE)
node(xs[5], y_main, nw, nh, "Validator", ORANGE)
node(xs[6], y_main, Inches(1.2), nh, "Report", NEUTRAL)

# Flèches principales
def cx(x, w=nw): return x + w
arrow(cx(xs[0]), y_main + nh//2, xs[1], y_main + nh//2)
arrow(cx(xs[1]), y_main + nh//2, xs[2], Inches(2.9))
arrow(cx(xs[1]), y_main + nh//2, xs[2], Inches(4.2))
arrow(cx(xs[2]), Inches(2.9), xs[3], y_main + nh//2)
arrow(cx(xs[2]), Inches(4.2), xs[3], y_main + nh//2)
arrow(cx(xs[3]), y_main + nh//2, xs[4], y_main + nh//2)
arrow(cx(xs[4]), y_main + nh//2, xs[5], y_main + nh//2)
arrow(cx(xs[5]), y_main + nh//2, xs[6], y_main + nh//2)

# Retry loop
loop = s.shapes.add_connector(1, cx(xs[5]) - Inches(0.2), y_main + nh,
                               xs[4] + Inches(0.4), y_main + nh)
loop.line.color.rgb = ORANGE; loop.line.width = Pt(1.5)
add_text(s, xs[4] + Inches(0.3), y_main + nh + Inches(0.05), Inches(1.5), Inches(0.3),
         "retry × 3", size=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Légende routage
y_leg = Inches(5.4)
legend = [
    ("Triage", "erreurs → Report directement", GREEN_DARK),
    ("Scanner", "0 finding → Report", GREEN_DARK),
    ("ExploitScorer", "rien d'exploitable → Report", PURPLE),
    ("Validator", "rejetés && iter<3 → Patcher", ORANGE),
]
col_w = Inches(3.0)
for i, (lbl, rule, col) in enumerate(legend):
    x = Inches(0.6) + i * (col_w + Inches(0.1))
    add_rect(s, x, y_leg, col_w, Inches(1.2), SLATE)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y_leg + Inches(0.3),
                              Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s, x + Inches(0.55), y_leg + Inches(0.25), col_w - Inches(0.7), Inches(0.35),
             lbl, size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), y_leg + Inches(0.65), col_w - Inches(0.4), Inches(0.5),
             rule, size=10, color=RGBColor(0xCB, 0xD5, 0xE1))

add_footer(s, 6)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Memory Safety Engine (Rust)
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "06 · MEMORY SAFETY", "Engine Rust externe",
            subtitle="20 règles couvrant buffer overflow, UAF, double-free, format string, unsafe Rust.")

# Bloc gauche : stats
left = Inches(0.6); top = Inches(2.4)
add_rect(s, left, top, Inches(5.0), Inches(4.2), DARK)
add_text(s, left + Inches(0.4), top + Inches(0.4), Inches(4), Inches(0.4),
         "COUVERTURE", size=11, bold=True, color=GREEN)
add_text(s, left + Inches(0.4), top + Inches(0.85), Inches(4), Inches(0.7),
         "20 règles · 3 langages", size=22, bold=True, color=WHITE)

stats = [("19", "Règles C/C++"), ("5", "Règles Rust unsafe"), ("8", "CWE distincts couverts")]
y = top + Inches(1.85)
for val, lbl in stats:
    add_text(s, left + Inches(0.4), y, Inches(1.5), Inches(0.6),
             val, size=28, bold=True, color=GREEN)
    add_text(s, left + Inches(1.8), y + Inches(0.15), Inches(3), Inches(0.5),
             lbl, size=13, color=RGBColor(0xCB, 0xD5, 0xE1))
    y += Inches(0.7)

# Bloc droit : exemples
right = Inches(5.8); top_r = Inches(2.4)
add_rect(s, right, top_r, Inches(6.93), Inches(4.2), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
add_text(s, right + Inches(0.4), top_r + Inches(0.4), Inches(6), Inches(0.4),
         "EXEMPLES DE RÈGLES", size=11, bold=True, color=GREEN)

rules = [
    ("MEM-001", "CWE-121", "strcpy() — stack overflow", "Critical"),
    ("MEM-010", "CWE-416", "use-after-free", "Critical"),
    ("MEM-050", "CWE-134", "format string non contrôlée", "Critical"),
    ("MEM-040", "CWE-190", "integer overflow dans malloc(n*m)", "High"),
    ("MEM-070", "CWE-119", "unsafe {} block (Rust)", "High"),
    ("MEM-073", "CWE-416", "Box::from_raw (Rust)", "High"),
]
y = top_r + Inches(0.85)
for rid, cwe, desc, sev in rules:
    add_text(s, right + Inches(0.4), y, Inches(1.1), Inches(0.4),
             rid, size=11, bold=True, color=DARK, font="Consolas")
    add_text(s, right + Inches(1.55), y, Inches(1.1), Inches(0.4),
             cwe, size=11, color=BLUE, font="Consolas")
    add_text(s, right + Inches(2.7), y, Inches(3.2), Inches(0.4),
             desc, size=11, color=GREY_TEXT)
    sev_color = RGBColor(0xDC, 0x26, 0x26) if sev == "Critical" else ORANGE
    pill = add_rect(s, right + Inches(6.0), y + Inches(0.04), Inches(0.85), Inches(0.3),
                    sev_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, right + Inches(6.0), y + Inches(0.04), Inches(0.85), Inches(0.3),
             sev, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.5)

add_footer(s, 7)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — LLM (Groq + NVIDIA)
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "07 · LLM", "Deux modèles, deux usages",
            subtitle="Fast pour le volume, Strong pour le raisonnement profond.")

card_w = Inches(5.95); card_h = Inches(4.2); top = Inches(2.4)

# Fast — Groq
add_rect(s, Inches(0.6), top, card_w, card_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
add_text(s, Inches(0.95), top + Inches(0.35), Inches(3), Inches(0.4),
         "FAST · GROQ", size=11, bold=True, color=PURPLE)
add_text(s, Inches(0.95), top + Inches(0.75), Inches(5), Inches(0.6),
         "Llama-3.1-8B-instant", size=20, bold=True, color=DARK)
add_text(s, Inches(0.95), top + Inches(1.4), Inches(5), Inches(0.4),
         "Latence typique  ·  200–500 ms", size=12, color=GREY_TEXT)
add_text(s, Inches(0.95), top + Inches(1.95), Inches(5), Inches(0.4),
         "UTILISÉ PAR", size=10, bold=True, color=GREY_TEXT)
for i, item in enumerate(["ExploitScorerAgent  ·  CVSS 3.1 + exploitabilité",
                          "Tâches à fort volume, scoring rapide"]):
    add_text(s, Inches(0.95), top + Inches(2.3) + i * Inches(0.4), Inches(5.3), Inches(0.4),
             "•  " + item, size=12, color=DARK)

# Strong — NVIDIA
right = Inches(6.78)
add_rect(s, right, top, card_w, card_h, DARK)
add_text(s, right + Inches(0.35), top + Inches(0.35), Inches(3), Inches(0.4),
         "STRONG · NVIDIA", size=11, bold=True, color=GREEN)
add_text(s, right + Inches(0.35), top + Inches(0.75), Inches(5), Inches(0.6),
         "Llama-3.1-70B-instruct", size=20, bold=True, color=WHITE)
add_text(s, right + Inches(0.35), top + Inches(1.4), Inches(5), Inches(0.4),
         "Latence typique  ·  1–3 s", size=12, color=RGBColor(0xCB, 0xD5, 0xE1))
add_text(s, right + Inches(0.35), top + Inches(1.95), Inches(5), Inches(0.4),
         "UTILISÉ PAR", size=10, bold=True, color=GREEN)
for i, item in enumerate(["SemanticAnalystAgent  ·  failles logiques + RAG",
                          "PatcherAgent  ·  unified diffs corrigeant la vuln"]):
    add_text(s, right + Inches(0.35), top + Inches(2.3) + i * Inches(0.4), Inches(5.3), Inches(0.4),
             "▸  " + item, size=12, color=WHITE)

add_footer(s, 8)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Mémoire persistante (3 couches)
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "08 · MÉMOIRE", "Trois couches de persistance",
            subtitle="SQLite pour l'identité, Qdrant pour la sémantique, RAM pour la session.")

layers = [
    ("Couche 1", "SQLite", ".memory_cache.db", GREEN,
     ["users", "user_projects", "scan_history (JSON)"],
     "Identité utilisateurs, projets et\nhistorique des scans."),
    ("Couche 2", "Qdrant", "vector store · 384 dim", PURPLE,
     ["vuln_patterns (Sentence-BERT)", "patches par CWE"],
     "RAG pour SemanticAnalyst & Patcher.\nApprentissage des patches validés."),
    ("Couche 3", "Session RAM", "in-process", ORANGE,
     ["SessionMemory", "non persistante"],
     "Cache court terme dans le process\nFastAPI pendant un scan."),
]

card_w = Inches(3.95); card_h = Inches(4.3); gap = Inches(0.18)
total_w = card_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) // 2
y = Inches(2.4)

for i, (eyebrow, name, tech, color, items, desc) in enumerate(layers):
    x = start_x + i * (card_w + gap)
    add_rect(s, x, y, card_w, card_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    # bandeau
    add_rect(s, x, y, card_w, Inches(0.7), color, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.6), Inches(0.4),
             eyebrow.upper(), size=10, bold=True, color=WHITE)
    add_text(s, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(0.5),
             name, size=20, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), y + Inches(1.4), card_w - Inches(0.6), Inches(0.4),
             tech, size=11, color=GREY_TEXT, font="Consolas")
    # items
    add_text(s, x + Inches(0.3), y + Inches(1.95), card_w - Inches(0.6), Inches(0.3),
             "CONTENU", size=9, bold=True, color=GREY_TEXT)
    for j, it in enumerate(items):
        add_text(s, x + Inches(0.3), y + Inches(2.3) + j * Inches(0.35),
                 card_w - Inches(0.6), Inches(0.35),
                 "•  " + it, size=11, color=DARK)
    add_text(s, x + Inches(0.3), y + Inches(3.45), card_w - Inches(0.6), Inches(0.8),
             desc, size=11, color=GREY_TEXT)

add_footer(s, 9)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Stack technique
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, LIGHT_BG)
title_block(s, "09 · STACK", "Technologies utilisées",
            subtitle="Choix techniques et leur rôle dans le pipeline.")

stack = [
    ("FastAPI", "0.136", "API REST asynchrone avec\nbackground tasks", GREEN),
    ("LangGraph", "1.2", "Orchestration conditionnelle\ndu pipeline d'agents", PURPLE),
    ("Groq + NVIDIA", "—", "Llama-3.1 8B (fast)\net 70B (strong)", BLUE),
    ("Semgrep", "1.163", "Analyse statique multi-langage\nrègles OWASP Top 10", ORANGE),
    ("Qdrant", "1.18", "Vector store pour\nles patterns sémantiques", PURPLE),
    ("Rust", "stable", "Memory safety engine\n20 règles compilées natif", BLUE),
]

cols = 3
card_w = Inches(4.0); card_h = Inches(2.05); gap = Inches(0.15)
total_w = card_w * cols + gap * (cols - 1)
start_x = (SLIDE_W - total_w) // 2
start_y = Inches(2.4)
for i, (name, ver, desc, color) in enumerate(stack):
    r, c = divmod(i, cols)
    x = start_x + c * (card_w + gap)
    y = start_y + r * (card_h + gap)
    add_rect(s, x, y, card_w, card_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    # barre verticale couleur
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, x + Inches(0.35), y + Inches(0.3), card_w - Inches(0.6), Inches(0.5),
             name, size=17, bold=True, color=DARK)
    add_text(s, x + Inches(0.35), y + Inches(0.8), card_w - Inches(0.6), Inches(0.35),
             "v " + ver, size=10, color=color, bold=True, font="Consolas")
    add_text(s, x + Inches(0.35), y + Inches(1.2), card_w - Inches(0.6), Inches(0.8),
             desc, size=11, color=GREY_TEXT)

add_footer(s, 10)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Conclusion
# ────────────────────────────────────────────────────────────────────────────
s = add_slide(); set_bg(s, DARK)
halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-4), Inches(4), Inches(8), Inches(8))
halo.fill.solid(); halo.fill.fore_color.rgb = GREEN_DARK
halo.line.fill.background(); halo.shadow.inherit = False

add_text(s, Inches(0.8), Inches(0.8), Inches(3), Inches(0.4),
         "CONCLUSION", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.8), Inches(1.3), Inches(12), Inches(1.2),
         "Ce que MAS apporte", size=40, bold=True, color=WHITE)

points = [
    ("Détection plus large",
     "5 outils statiques + LLM 70B couvrent ce que Semgrep seul rate :\nfailles logiques, IDOR, race conditions."),
    ("Priorisation utile",
     "Chaque vuln reçoit un CVSS 3.1 + un booléen d'exploitabilité.\nFini la liste plate de 200 warnings."),
    ("Patches automatisés",
     "Diff généré, appliqué, re-scanné, retry × 3 si rejeté.\nLes patches validés enrichissent Qdrant."),
    ("Mémoire qui apprend",
     "Patterns et patches s'accumulent dans Qdrant —\nla qualité augmente scan après scan."),
]
y = Inches(2.9)
for title, desc in points:
    add_rect(s, Inches(0.8), y, Inches(0.08), Inches(0.85), GREEN, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(1.05), y + Inches(0.02), Inches(11), Inches(0.4),
             title, size=15, bold=True, color=WHITE)
    add_text(s, Inches(1.05), y + Inches(0.42), Inches(11), Inches(0.5),
             desc, size=11, color=RGBColor(0xCB, 0xD5, 0xE1))
    y += Inches(1.0)

add_text(s, Inches(0.8), Inches(6.95), Inches(11), Inches(0.4),
         "Documentation complète  ·  /docs  ·  index, architecture, agents, workflow, memory, llm, api, ui",
         size=10, color=GREY_TEXT)
add_text(s, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3),
         "11 / 11", size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ────────────────────────────────────────────────────────────────────────────
prs.save(str(OUTPUT))
print(f"OK -> {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
